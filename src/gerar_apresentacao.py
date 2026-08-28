from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .config import OUTPUT_DIR

AZUL = RGBColor(0x10, 0x4F, 0x8A)
AZUL_CLARO = RGBColor(0x2C, 0x7F, 0xB8)
VERDE = RGBColor(0x1E, 0x7E, 0x34)
CINZA = RGBColor(0x55, 0x55, 0x55)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
COR_LADO = RGBColor(0xE8, 0xF1, 0xF8)
PRETO = RGBColor(0x22, 0x22, 0x22)


def nova_slide(prs, titulo: str, subtitulo: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # barra lateral cor
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL
    bar.line.fill.background()
    # faixa título
    faixa = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.25), Inches(9.5), Inches(0.9))
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = COR_LADO
    faixa.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(9.0), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = AZUL
    if subtitulo:
        p2 = tf.add_paragraph()
        p2.text = subtitulo
        p2.font.size = Pt(13)
        p2.font.color.rgb = CINZA
        p2.font.italic = True
    return slide


def add_texto(slide, x, y, w, h, linhas, tamanho=16, negrito=False, cor=PRETO, bullet=False, espaco=8):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, linha in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if bullet else "") + linha
        p.font.size = Pt(tamanho)
        p.font.bold = negrito
        p.font.color.rgb = cor
        p.space_after = Pt(espaco)
    return tb


def add_imagem(slide, path: Path, x, y, w=None, h=None):
    kwargs = {"left": Inches(x), "top": Inches(y)}
    if w:
        kwargs["width"] = Inches(w)
    if h:
        kwargs["height"] = Inches(h)
    slide.shapes.add_picture(str(path), **kwargs)


def design_titulo(slide):
    # cantinho "jt2026 · Seazone · Itapema/SC"
    add_texto(slide, 0.8, 6.9, 9.0, 0.4, ["Seazone · Itapema/SC · Jovens Talentos 2026"], tamanho=10, cor=CINZA)


def gerar():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ------- SLIDE 1 — Capa -------
    s = nova_slide(prs, "")
    add_texto(s, 0.8, 2.0, 8.5, 0.8, ["Recomendação de Investimento — Short Stay", "Itapema/SC"], tamanho=32, negrito=True, cor=AZUL)
    add_texto(s, 0.8, 3.2, 8.5, 0.5, ["Para a Seazone · Análise apoiada em IA (Fases 0–7 + ai-log)"], tamanho=16, cor=CINZA)
    add_texto(s, 0.8, 4.0, 8.5, 0.5, ["Vídeo: 3 min · Estrutura: recomendaão → racional → processo → próximos passos"], tamanho=13, cor=CINZA)

    # ------- SLIDE 2 — Recomendação (resultado primeiro) -------
    s = nova_slide(prs, "1 · O que fazer — em 2 frases", "resultado primeiro")
    add_texto(s, 1.0, 1.5, 8.3, 0.6, ["COMPRAR: apartamentos compactos de 1 quarto"], tamanho=24, negrito=True, cor=AZUL)
    add_texto(s, 1.0, 2.2, 8.3, 0.5, ["Localização: Morretes e Meia Praia — NÃO o Centro"], tamanho=20, negrito=True, cor=VERDE)
    add_texto(s, 1.0, 3.0, 8.3, 0.4, ["Execução: 60% originação/lançamento + 40% pronto (piloto)"], tamanho=16, bullet=True)
    add_texto(s, 1.0, 3.5, 8.3, 0.4, ["Morretes, 55m² construído: yield 2,25% (conservador) → 7,1% (otimista, occ 53%)"], tamanho=16, bullet=True)
    design_titulo(s)

    # ------- SLIDE 3 — Critério -------
    s = nova_slide(prs, "2 · O critério: o que é “melhor”?", "racional da decisão")
    add_texto(s, 1.0, 1.6, 8.0, 0.6, ["'Melhor' = maior YIELD LÍQUIDO anual (NOI ÷ investimento total)"], tamanho=18, negrito=True, cor=AZUL)
    add_texto(s, 1.0, 2.3, 8.0, 0.5, ["Não é a maior receita bruta — é eficiência de capital no longo prazo"], tamanho=14, bullet=True)
    add_texto(s, 1.0, 3.0, 8.0, 0.5, ["Três eixos de evidência (com números):"], tamanho=14, bullet=True)
    add_texto(s, 1.0, 3.6, 8.0, 0.5, ["Perfil: +1 quarto (mesmos hóspedes) ⇒ receita −29% | +1 hóspede ⇒ +34% | reviews dobram ⇒ +39%"], tamanho=14, bullet=True)
    add_texto(s, 1.0, 4.1, 8.0, 0.5, ["Localização: m² Morretes R$11.6k vs Centro R$16.8k (44% mais caro) — quem compra barato ganha"], tamanho=14, bullet=True)
    add_texto(s, 1.0, 4.6, 8.0, 0.5, ["Operação: reserva instantânea ⇒ +105% receita"], tamanho=14, bullet=True)
    design_titulo(s)

    # ------- SLIDE 4 — Evidência 1: perfil (receita por quartos) + coef -------
    s = nova_slide(prs, "Evidência 1 — Perfil compacto vence", "Fase 3 + Fase 4")
    add_imagem(s, OUTPUT_DIR / "fase3_barra_quartos.png", 0.7, 1.3, w=4.6)
    add_imagem(s, OUTPUT_DIR / "fase4_coef_plot.png", 5.3, 1.3, w=4.4)
    add_texto(s, 0.8, 6.6, 8.5, 0.5, ["Receita cresce com quartos, mas custo/compra cresce mais — o yield é decidido pela eficiência de capital"], tamanho=12, cor=CINZA)
    design_titulo(s)

    # ------- SLIDE 5 — Evidência 2: localização (boxplot + ranking) -------
    s = nova_slide(prs, "Evidência 2 — Onde o dinheiro rende", "Fase 3")
    add_imagem(s, OUTPUT_DIR / "fase3_ranking_bairro.png", 0.7, 1.3, w=5.2)
    add_imagem(s, OUTPUT_DIR / "fase3_boxplot_bairro.png", 5.9, 1.3, w=3.8)
    add_texto(s, 0.8, 6.6, 8.5, 0.5, ["Mediana de receita: Meia Praia R$3,1k/mês (n=632) · Centro R$2,3k (n=205) · Morretes R$2,1k (n=83) — mas m² muda o jogo"], tamanho=12, cor=CINZA)
    design_titulo(s)

    # ------- SLIDE 6 — Tese dos compactos (veredito) -------
    s = nova_slide(prs, "3 · Veredito: tese dos compactos no Centro", "posição clara")
    add_texto(s, 1.0, 1.5, 8.0, 0.6, ["SUSTENTA PARCIALMENTE — com correção de bairro"], tamanho=20, negrito=True, cor=VERDE)
    add_texto(s, 1.0, 2.2, 8.0, 0.5, ["✔ Perfil compacto: supera imóveis maiores em yield em todos os cenários"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 2.7, 8.0, 0.5, ["✗ Localização: o melhor bairro para compactos é Morretes (+3,4% otimista), não o Centro (+0,6%)"], tamanho=15, bullet=True)
    add_imagem(s, OUTPUT_DIR / "fase5_veredito.png", 1.0, 3.6, w=8.0)
    add_texto(s, 1.0, 6.8, 8.0, 0.4, ["A tese acertou no 'o quê', errou no 'onde'"], tamanho=13, cor=CINZA)
    design_titulo(s)

    # ------- SLIDE 7 — Estimativa de retorno -------
    s = nova_slide(prs, "O que comprar e o retorno", "Morretes · 1 quarto · 55m²")
    linhasA = ["Invest: R$727k", "NOI: R$2,4k/ano (base)", "Yield: 0,33% → 4,0% (otim)", "occ 12% ⇒ negativo"]
    linhasB = ["Invest: R$648k", "NOI: R$14,6k/ano (base)*", "Yield: 2,25% → 7,1% (otim)*", "occ 13% ⇒ negativo"]

    add_texto(s, 1.0, 1.4, 8.0, 0.5, ["Unidade compacta 1q — execução híbrida 60/40"], tamanho=18, negrito=True, cor=AZUL)
    # coluna A
    add_texto(s, 1.0, 2.0, 3.8, 2.0, ["PRONTO (A)"], tamanho=15, negrito=True, cor=AZUL)
    add_texto(s, 1.0, 2.5, 3.8, 2.5, linhasA, tamanho=13, bullet=True, espaco=6)
    # coluna B
    add_texto(s, 5.2, 2.0, 4.3, 2.0, ["LANÇAMENTO (B)"], tamanho=15, negrito=True, cor=VERDE)
    add_texto(s, 5.2, 2.5, 4.3, 2.5, linhasB, tamanho=13, bullet=True, espaco=6)
    add_imagem(s, OUTPUT_DIR / "fase6_cumulative_noi.png", 3.5, 4.5, w=6.0)
    add_texto(s, 1.0, 6.95, 8.0, 0.4, ["*Lançamento inclui prêmio de ativo novo (+8% diária, +10% occ) e produção 75% da revenda"], tamanho=11, cor=CINZA)
    design_titulo(s)

    # ------- SLIDE 8 — Como a IA foi usada -------
    s = nova_slide(prs, "4 · Como usei IA (e onde a critiquei)", "processo avaliado")
    add_texto(s, 1.0, 1.6, 8.2, 0.5, ["IA em todas as fases · conversa inteira exportada em ai-log/ (170 mensagens)"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 2.2, 8.2, 0.5, ["Obstáculo real (bug de processo): 1ª régua comparava construir R$4,2k/m² com comprar R$16k/m² — "], tamanho=15, bullet=True)
    add_texto(s, 1.0, 2.6, 8.2, 0.5, ["  diferença de 4x inexistente no mercado"], tamanho=13, bullet=True, cor=CINZA)
    add_texto(s, 1.0, 3.3, 8.2, 0.5, ["Como driblei: refiz a régua com produção ≈ 75% da revenda e re-testei a tese"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 4.0, 8.2, 0.5, ["Veredito honesto (contra o primeiro 'sustenta'): a tese acerta no perfil, erra no bairro"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 4.7, 8.2, 0.5, ["Lição: IA acelera iteração; a decisão final passou pelo meu crivo"], tamanho=15, bullet=True, cor=VERDE)
    design_titulo(s)

    # ------- SLIDE 9 — Próximos passos (+1 semana) -------
    s = nova_slide(prs, "O que faria com +1 semana", "fechamento")
    add_texto(s, 1.0, 1.6, 8.0, 0.5, ["1. Validar ocupação com calendário real (proxy captura = hoje)"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 2.2, 8.0, 0.5, ["2. Fechar orçamento de obra e VGV do lançamento (produção real)"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 2.8, 8.0, 0.5, ["3. Sazonalidade de alta temporada (jan/fev) para calibrar cenários"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 3.4, 8.0, 0.5, ["4. Amenidades que mais convertem, controlando o tamanho do imóvel"], tamanho=15, bullet=True)
    add_texto(s, 1.0, 4.4, 8.2, 0.6, ["Núcleo inalterado: compacto, fora do Centro, com ocupação ≥30% via gestão de canal"], tamanho=16, negrito=True, cor=AZUL)
    add_texto(s, 1.0, 5.2, 8.2, 0.5, ["Sem 30% de ocupação, nenhum caminho fecha."], tamanho=14, cor=CINZA, bullet=True)
    # close
    add_texto(s, 1.0, 6.2, 8.0, 0.5, ["Análise completa + processo: repo jt2026-caio-oliveira"], tamanho=12, cor=CINZA)
    design_titulo(s)

    destino = OUTPUT_DIR / "apresentacao_apoio_video.pptx"
    prs.save(str(destino))
    print("Gerado:", destino)
    print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))


if __name__ == "__main__":
    gerar()